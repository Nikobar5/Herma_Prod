import os
import openpyxl
import fitz
import aiofiles
import chardet

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.schema.document import Document
from get_embedding_function import get_embedding_function
from langchain_chroma import Chroma
from concurrent.futures import ThreadPoolExecutor
from docx import Document as DocxDocument
from pptx import Presentation
from pathlib import Path
import time

class Uploaded_data:
    def __init__(self, name, data_path, non_chat_history, chunk_size):
        self.non_chat_history = non_chat_history
        self.name = name
        self.chunk_size = chunk_size
        self.data_path = data_path
        self.documents = self.load_documents(self.data_path)
        self.timestamp = int(time.time() * 1000)
        self.vector_database_path = f"{name}_{self.timestamp}"

        self.add_to_chroma()

        if non_chat_history:
            self.data_summary = self.generate_summary()




    def load_documents(self, data_path):
        if not os.path.exists(data_path):
            raise FileNotFoundError(f"File not found: {data_path}")
        if os.path.isdir(data_path):
            raise FileNotFoundError("Directory loading not supported. Provide a single file.")
        elif data_path.lower().endswith(".pdf"):
            return self._process_pdf(data_path)
        elif data_path.lower().endswith((".txt", ".md")):
            return self._process_text(data_path)
        elif data_path.lower().endswith(".docx"):
            return self._process_word(data_path)
        elif data_path.lower().endswith(".pptx"):
            return self._process_pptx(data_path)
        elif data_path.lower().endswith(".xlsx"):
            return self._process_excel(data_path)
        elif data_path.lower().endswith(".csv"):
            return self._process_csv(data_path)
        elif data_path.lower().endswith(".json"):
            return self._process_json(data_path)
        else:
            raise ValueError(f"Unsupported file type for: {data_path}")
        try:
            return document_loader.load()
        except Exception as e:
            raise RuntimeError(f"Failed to load documents from {data_path}: {e}")

    def _process_pdf(self, pdf_path):
        documents = []
        doc = fitz.open(pdf_path)

        for i, page in enumerate(doc):
            text = page.get_text("text")
            table_text = self._extract_tables_from_pdf(page)
            combined_text = f"{text}\n\nTables:\n{table_text}" if table_text else text

            documents.append(Document(page_content=combined_text, metadata={"source": pdf_path, "page": i + 1}))

        return documents

    def _extract_tables_from_pdf(self, page):
        markdown_tables = []
        raw_text = page.get_text("text")
        if not raw_text:
            return ""

        rows = raw_text.split("\n")

        table_rows = []
        in_table = False

        for row in rows:
            if "\t" in row or "  " in row:
                if not in_table:
                    in_table = True
                    table_rows = []


                cells = []
                if "\t" in row:
                    cells = row.split("\t")
                else:

                    import re
                    cells = re.split(r'  +', row)


                cells = [cell.strip() for cell in cells]
                table_rows.append(cells)
            else:

                if in_table and table_rows:
                    markdown_table = self._convert_to_markdown_table(table_rows)
                    markdown_tables.append(markdown_table)
                    in_table = False


        if in_table and table_rows:
            markdown_table = self._convert_to_markdown_table(table_rows)
            markdown_tables.append(markdown_table)

        return "\n\n".join(markdown_tables)

    def _extract_images_from_pdf(self, page):
        extracted_texts = []
        image_list = page.get_images(full=True)

        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            img_path = f"temp_image_{xref}.png"
            base_image.pil_save(img_path)

            img_text = self._process_image(img_path)
            os.remove(img_path)

            extracted_texts.append(
                (img_text[0].page_content, {"source": page.parent.name, "page": page.number + 1, "type": "image"}))

        return extracted_texts

    def _convert_to_markdown_table(self, table_rows):

        if not table_rows or len(table_rows) == 0:
            return ""

        max_cols = max(len(row) for row in table_rows)

        normalized_rows = []
        for row in table_rows:
            if len(row) < max_cols:
                normalized_rows.append(row + [''] * (max_cols - len(row)))
            else:
                normalized_rows.append(row)

        markdown_lines = []

        header = normalized_rows[0]
        markdown_lines.append('| ' + ' | '.join(header) + ' |')

        markdown_lines.append('| ' + ' | '.join(['---'] * len(header)) + ' |')

        for row in normalized_rows[1:]:
            markdown_lines.append('| ' + ' | '.join(row) + ' |')

        return '\n'.join(markdown_lines)

    def split_documents(self):
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=50,
            length_function=len,
            is_separator_regex=False,
        )
        with ThreadPoolExecutor() as executor:
            results = list(executor.map(text_splitter.split_documents, [self.documents]))
            return [chunk for sublist in results for chunk in sublist]

    def _process_text(self, text_path):
        with open(text_path, "r", encoding="utf-8") as f:
            text = f.read()
        return [Document(page_content=text, metadata={"source": text_path})]

    def _process_word(self, word_path):
        try:
            doc = DocxDocument(word_path)
            text = "\n".join([p.text for p in doc.paragraphs])
            return [Document(page_content=text, metadata={"source": word_path})]
        except Exception as e:
            raise RuntimeError(f"Failed to process Word document {word_path}: {e}")

    def _process_pptx(self, pptx_path):
        try:
            prs = Presentation(pptx_path)
            text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            return [Document(page_content=text, metadata={"source": pptx_path})]
        except Exception as e:
            raise RuntimeError(f"Failed to process PowerPoint file {pptx_path}: {e}")

    def _process_excel(self, excel_path):
        try:
            wb = openpyxl.load_workbook(excel_path)
            text = "\n".join(
                ["\t".join([str(cell.value) for cell in row]) for sheet in wb for row in sheet.iter_rows()])
            return [Document(page_content=text, metadata={"source": excel_path})]
        except Exception as e:
            raise RuntimeError(f"Failed to process Excel file {excel_path}: {e}")

    async def _process_csv(self, csv_path):
        with open(csv_path, "rb") as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result["encoding"]
        if encoding != "utf-8":
            encoding = "utf-8"
        async with aiofiles.open(csv_path, "r", encoding=encoding, errors="replace") as f:
            text = await f.read()
        return [Document(page_content=text, metadata={"source": csv_path})]

    async def _process_json(self, json_path):
        async with aiofiles.open(json_path, "r", encoding="utf-8") as f:
            text = await f.read()
        return [Document(page_content=text, metadata={"source": json_path})]

    def add_to_chroma(self):
        chunks = self.split_documents()
        db_path = self.get_db_path()
        db = Chroma(
            persist_directory=str(db_path),
            embedding_function=get_embedding_function()
        )

        chunks_with_ids = self.calculate_chunk_ids(chunks)

        existing_items = db.get(include=[])
        existing_ids = set(existing_items["ids"])
        new_chunks = chunks_with_ids
        new_chunk_ids = [chunk.metadata["id"] for chunk in new_chunks]
        db.add_documents(new_chunks, ids=new_chunk_ids)
        end_time = time.time()

    def calculate_chunk_ids(self, chunks):
        last_page_id = None
        current_chunk_index = 0
        for chunk in chunks:
            source = chunk.metadata.get("source", "unknown")
            page = chunk.metadata.get("page", 0)
            current_page_id = f"{source} Page: {page}"
            if current_page_id == last_page_id:
                current_chunk_index += 1
            else:
                current_chunk_index = 0
            chunk_id = f"{current_page_id}:{current_chunk_index}"
            last_page_id = current_page_id
            chunk.metadata["id"] = chunk_id

        return chunks

    def generate_summary(self):
        from langchain_ollama import ChatOllama
        chunks = self.split_documents()
        if not chunks:
            return "No content available for summarization."
        sample_chunks = []
        is_full_document = len(chunks) < 6
        if not is_full_document:
            sample_chunks = chunks[:3] + chunks[-3:]

        else:
            sample_chunks = chunks
        sample_text = "\n\n---\n\n".join([chunk.page_content for chunk in sample_chunks])
        if is_full_document:
            summary_prompt = f"""
            Below is the full text of a document titled '{self.name}'. 

            Please provide a TWO-SENTENCE summary of what this document is about.
            Include any information about: title, author, publishing date, main topic, key implications. 
            Format your response without bullet points.

            DOCUMENT TEXT:
            {sample_text}
            """
        else:
            middle_index = len(sample_text) // 2
            first_half = sample_text[:middle_index]
            second_half = sample_text[middle_index:]
            summary_prompt = f"""
            Below is text from the first and last parts of a document titled '{self.name}'. 

            Please provide a TWO-SENTENCE summary of what this document is about.
            Include any information about: title, author, publishing date, main topic, key implications. 
            Format your response without bullet points.

            DOCUMENT TEXT:
            Here are the first two chunks at the beginning, likely containing title, author, date, intro
            {first_half}
            Here are the last two chunks at the end of the doc, likely containing conclusion or references
            {second_half}
            """
        llm = ChatOllama(model="llama3.2:1b", temperature=0.5, num_predict=100)
        result = llm.invoke(summary_prompt)

        return result.content

    @staticmethod
    def get_project_root():
        return Path(__file__).resolve().parents[2]

    def get_db_path(self):
        return self.get_project_root() / 'storage' / 'db_store' / self.vector_database_path

    @staticmethod
    def delete_vector_db(filename):
        db_root = Uploaded_data.get_project_root() / 'storage' / 'db_store'

        matching_dbs = [d for d in os.listdir(str(db_root))
                        if d.startswith(filename + "_") or d == filename]

        for db_name in matching_dbs:
            db_path = db_root / db_name
            if os.path.exists(str(db_path)):
                try:
                    import shutil
                    shutil.rmtree(str(db_path))
                except Exception as e:
                    print(f"Error deleting database {db_path}: {e}")