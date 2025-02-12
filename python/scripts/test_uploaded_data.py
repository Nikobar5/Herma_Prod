import os
import shutil
import json
import pandas as pd
import openpyxl
from PIL import Image, ImageDraw
from docx import Document as DocxDocument
from pptx import Presentation
from uploaded_data import Uploaded_data

import subprocess

subprocess.check_call(['pip', 'install', 'fpdf'])
from fpdf import FPDF

# Create test files in a temporary directory
TEMP_DIR = "temp_test_files"


def setup_test_environment():
    """Creates test files for different formats."""
    os.makedirs(TEMP_DIR, exist_ok=True)

    # PDF File (Placeholder)
    os.makedirs(TEMP_DIR, exist_ok=True)
    create_pdf(os.path.join(TEMP_DIR, "test.pdf"))

    # Image File
    img = Image.new('RGB', (100, 100), color=(73, 109, 137))
    draw = ImageDraw.Draw(img)
    draw.text((10, 10), "Test Image", fill=(255, 255, 255))
    img.save(os.path.join(TEMP_DIR, "test.png"))

    # Text File
    with open(os.path.join(TEMP_DIR, "test.txt"), "w") as f:
        f.write("This is a test text file.")

    # Word Document
    doc = DocxDocument()
    doc.add_paragraph("This is a test Word document.")
    doc.save(os.path.join(TEMP_DIR, "test.docx"))

    # PowerPoint File
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    textbox = slide.shapes.add_textbox(10, 10, 500, 100)
    textbox.text = "Test PowerPoint Slide"
    ppt.save(os.path.join(TEMP_DIR, "test.pptx"))

    # Excel File
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Column1", "Column2"])
    ws.append(["Data1", "Data2"])
    wb.save(os.path.join(TEMP_DIR, "test.xlsx"))

    # CSV File
    df = pd.DataFrame({"Column1": ["Data1", "Data2"], "Column2": ["More1", "More2"]})
    df.to_csv(os.path.join(TEMP_DIR, "test.csv"), index=False)

    # JSON File
    json_data = {"key1": "value1", "key2": "value2"}
    with open(os.path.join(TEMP_DIR, "test.json"), "w") as f:
        json.dump(json_data, f, indent=4)


def create_pdf(file_path):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="This is a test PDF file.", ln=True, align="C")
    pdf.output(file_path)


def test_uploaded_data():
    """Runs tests on all supported file types."""
    test_files = [
        "test.pdf", "test.png", "test.txt", "test.docx",
        "test.pptx", "test.xlsx", "test.csv", "test.json"
    ]

    for file_name in test_files:
        file_path = os.path.join(TEMP_DIR, file_name)
        print(f"\nTesting file: {file_path}")

        try:
            uploaded_data = Uploaded_data(name=file_name, data_path=file_path)
            assert uploaded_data.documents, "No documents loaded"
            print(f"✅ {file_name}: Successfully processed.")

        except Exception as e:
            print(f"❌ {file_name}: Failed with error - {e}")


def cleanup_test_environment():
    """Deletes the temporary test files and directories."""
    shutil.rmtree(TEMP_DIR, ignore_errors=True)
    print("\n🧹 Test environment cleaned up.")


if __name__ == "__main__":
    setup_test_environment()
    test_uploaded_data()
    cleanup_test_environment()